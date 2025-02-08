
import pdfplumber
import re
from bs4 import BeautifulSoup

import uuid

from resources import db

from firebase_admin import  firestore
import os
from datetime import datetime, timezone
from typing import List
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib import utils
from pptx import Presentation
from pptx.util import Inches, Pt

from reportlab.platypus import Table, TableStyle

from reportlab.lib import colors


LOCAL_PATH = os.getenv("LOCAL_PATH")

def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text




def is_valid_url(url):
    regex = re.compile(
        r"^(http|https)://"  # http:// or https://
        r"([A-Za-z0-9.-]+)"  # domain...
        r"(:\d+)?"  # optional port
        r"(/.*)?$"  # optional path
    )
    return re.match(regex, url) is not None


def extract_text_from_website(response_text:str):
    soup = BeautifulSoup(response_text, "html.parser")
    url_content = soup.get_text(separator=" ", strip=True)
    return url_content


def upload_content(content, type, filename=None):
    uuid_ = str(uuid.uuid4())
    doc_ref = db.collection("items").document(uuid_)
    doc_ref.set({"content": content, "type": type,"filename":filename})
    return uuid_


async def write_file_content(filepath,filecontent):
    with open(filepath, "wb") as f:
        content= await filecontent.read()
        f.write(content)


def get_content_from_firebase(id, collection="items"):
    try:
        doc_ref = db.collection(collection).document(id)
        doc = doc_ref.get()
        data = doc.to_dict()
        return data
    except Exception as e:
        return None


def store_prompt(prompt, answer):
    current_utc_time_question = datetime.now(timezone.utc)
    uuid_ques = str(uuid.uuid4())
    doc_ref = db.collection("prompts").document(uuid_ques)
    doc_ref.set(
            {
                            "prompt": prompt,
                            "answer": answer,
                            "date_created": str(current_utc_time_question),
            }
    )


def get_prompts():
    collection_ref = db.collection("prompts")
    docs = collection_ref.order_by("date_created", direction=firestore.Query.DESCENDING).limit(5).get()
    documents = []
    for doc in docs:
        documents.append({**doc.to_dict(), "id": doc.id})  # Include document ID
    return documents


def find_prompts(ids: List[str]):
    collection_ref = db.collection("prompts")
    documents = []

    for doc_id in ids:
        doc_ref = collection_ref.document(doc_id)  # Reference to the document
        doc = doc_ref.get()  # Fetch the document
        if doc.exists:
            documents.append({**doc.to_dict(), "id": doc.id})  # Convert document to dictionary
    
    return documents


def wrap_text(text, width, font_size):
    """
    Wrap text to fit within the specified width.
    
    Args:
        text (str): The text to wrap.
        width (int): The width in points (1 point = 1/72 inch).
        font_size (int): The font size being used.

    Returns:
        list: List of wrapped lines.
    """
    words = text.split()
    lines = []
    current_line = ""
    max_chars_per_line = int(width / (font_size * 0.5))  # Approximate characters per line

    for word in words:
        if len(current_line) + len(word) + 1 <= max_chars_per_line:
            current_line += (word + " ")
        else:
            lines.append(current_line.strip())
            current_line = word + " "
    lines.append(current_line.strip())
    return lines

def generate_pdf_with_wrapping(data, layout, heading, subheading):
    file_path = os.path.join("outputs","output.pdf")
    c = canvas.Canvas(file_path, pagesize=A4)
    width, height = A4


    def draw_header():
        """Draw the header with the current date and 'SYSTEM GENERATED' text."""
        c.setFont("Helvetica", 10)
        current_date = datetime.now().strftime("%Y-%m-%d")
        c.drawString(50, height - 30, f"Date: {current_date}")
        c.drawRightString(width - 50, height - 30, "SYSTEM GENERATED")


    def render_table(c, table_data, x, y, max_width, font_size, page_height, line_spacing):
   
    # Wrap each cell's content to fit in the column width
        wrapped_data = []
        for row in table_data:
            wrapped_row = []
            for cell in row:
                wrapped_lines = wrap_text(cell, max_width / len(row), font_size)
                wrapped_row.append("\n".join(wrapped_lines))
            wrapped_data.append(wrapped_row)

        # Calculate the total height required for the table
        table_height = len(table_data) * (line_spacing * 2)  # Approximate height for rows with padding

        # Check if the table will fit on the current page
        if y - table_height < 50:
            c.showPage()  # Move to the next page
            y = page_height - 50  # Reset Y position for a new page

        # Create and draw the table
        col_widths = [max_width / len(table_data[0])] * len(table_data[0])
        table = Table(wrapped_data, colWidths=col_widths)
        style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ])
        table.setStyle(style)
        table.wrapOn(c, max_width, y)
        table.drawOn(c, x, y - table._height)
        return y - table._height - 20


    # First page
    draw_header()
    
    # Add heading and subheading
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(width / 2, height - 50, heading)
    c.setFont("Helvetica", 14)
    c.drawCentredString(width / 2, height - 80, subheading)
    
    y = height - 120
    font_size = 12
    line_spacing = 16  # Line spacing
    
    for i, doc in enumerate(data):
        if layout == "1perpage" or i % 2 == 0:
            if i != 0:  # Add a new page for subsequent data
                c.showPage()
                y = height - 50  # Reset y for a new page
                # Add heading/subheading on new page
                c.setFont("Helvetica-Bold", 20)
                c.drawCentredString(width / 2, y, heading)
                c.setFont("Helvetica", 14)
                c.drawCentredString(width / 2, y - 30, subheading)
                y -= 50
        
        # Add document content
        c.setFont("Helvetica", 12)
        prompt_lines = wrap_text(f"Prompt: {doc['prompt']}", width - 100, font_size)
        #answer_lines = wrap_text(f"Answer: {doc.get('answer', '')}", width - 100, font_size)

        for line in prompt_lines:
            c.drawString(50, y, line)
            y -= line_spacing
            if y < 50:  # Check if there's space on the page
                c.showPage()
                y = height - 50
        
        y -= 10  # Add some spacing between prompt and answer


        '''
        for line in answer_lines:
            c.drawString(50, y, line)
            y -= line_spacing
            if y < 50:  # Check if there's space on the page
                c.showPage()
                y = height - 50
'''

        answer = doc.get('answer', '')
        blocks = answer.split('\n')
        for block in blocks:
            block = block.strip()
            if block.startswith('|') and block.endswith('|'):  # Detect table block
                table_data = [row.strip('|').split('|') for row in block.split('\n') if row.startswith('|')]
                y = render_table(c, table_data, 50, y, width - 100, font_size, height, line_spacing)
            else:  # Render text block
                text_lines = wrap_text(block, width - 100, font_size)
                for line in text_lines:
                    c.drawString(50, y, line)
                    y -= line_spacing
                    if y < 50:
                        c.showPage()
                        
                        y = height - 120



        if layout == "2perpage" and i % 2 == 0:
            y = height / 2  # Adjust position for 2-per-page layout

    c.save()
    return file_path
# Utility: Generate PDF

'''

def generate_pdf(data, layout, heading, subheading):
    file_path = "/tmp/output.pdf"
    c = canvas.Canvas(file_path, pagesize=A4)
    width, height = A4
    
    # Add heading and subheading
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(width / 2, height - 50, heading)
    c.setFont("Helvetica", 14)
    c.drawCentredString(width / 2, height - 80, subheading)
    
    y = height - 120
    for i, doc in enumerate(data):
        if layout == "1perpage" or i % 2 == 0:
            if i != 0:  # Add a new page for subsequent data
                c.showPage()
                y = height - 50  # Reset y for a new page
            c.setFont("Helvetica-Bold", 20)
            c.drawCentredString(width / 2, y, heading)
            c.setFont("Helvetica", 14)
            c.drawCentredString(width / 2, y - 30, subheading)
            y -= 50  # Adjust position for the first content
        
        # Add document content
        c.setFont("Helvetica", 12)
        c.drawString(50, y, f"Prompt: {doc['prompt']}")
        y -= 20
        c.drawString(50, y, f"Answer: {doc.get('answer')}")
        y -= 40
        
        if layout == "2perpage" and i % 2 == 0:
            y = height / 2  # Adjust position for 2-per-page layout

    c.save()
    return file_path
'''


# Utility: Generate PPT
def generate_ppt(data, layout, heading, subheading):
    ppt = Presentation()
    
    for i, doc in enumerate(data):
        if layout == "1perpage" or i % 2 == 0:
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])
            # Add heading and subheading
            textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = textbox.text_frame
            tf.text = heading
            tf.paragraphs[0].font.size = Pt(24)
            p = tf.add_paragraph()
            p.text = subheading
            p.font.size = Pt(18)
        
        # Add document content
        y_offset = Inches(2 if layout == "1perpage" else 3) if i % 2 == 0 else Inches(4.5)
        content_box = slide.shapes.add_textbox(Inches(1), y_offset, Inches(8), Inches(2))
        content_tf = content_box.text_frame
        content_tf.text = f"Prompt: {doc['prompt']}\nAnswer: {doc.get('answer', '')}"
    
    file_path =os.path.join("outputs","output.pptx")
    ppt.save(file_path)
    return file_path
